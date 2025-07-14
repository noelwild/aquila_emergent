"""Document processing service."""

import hashlib
import base64
import aiofiles
from typing import List, Dict, Any, Callable
from pathlib import Path
import shutil
import os
from PyPDF2 import PdfReader
from pptx import Presentation
from openpyxl import load_workbook
from PIL import Image
import pytesseract
from docx import Document
from pdf2image import convert_from_path
from jinja2 import Environment, FileSystemLoader, select_autoescape
import xmlschema
from reportlab.platypus import SimpleDocTemplate, Paragraph, Image, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
import io
import asyncio
import uuid
import re
import logging
from datetime import datetime

from ..models.document import (
    UploadedDocument,
    ICN,
    DataModule,
    ProcessingTask,
    PublicationModule,
)
from ..models.base import DMTypeEnum, SettingsModel, StructureType, SecurityLevel
from ..ai_providers.provider_factory import ProviderFactory
from ..ai_providers.base import TextProcessingRequest, VisionProcessingRequest
from ..services.audit import AuditService

logger = logging.getLogger(__name__)

# Default structural codes based on S1000D operational environment
DEFAULT_STRUCTURE_CODES: Dict[StructureType, Dict[str, str]] = {
    StructureType.AIR: {
        "system_diff": "10",
        "system_code": "211",
        "sub_system_code": "00",
        "sub_sub_system_code": "00",
    },
    StructureType.WATER: {
        "system_diff": "20",
        "system_code": "311",
        "sub_system_code": "00",
        "sub_sub_system_code": "00",
    },
    StructureType.LAND: {
        "system_diff": "30",
        "system_code": "411",
        "sub_system_code": "00",
        "sub_sub_system_code": "00",
    },
    StructureType.OTHER: {
        "system_diff": "00",
        "system_code": "000",
        "sub_system_code": "00",
        "sub_sub_system_code": "00",
    },
}

# Info code mapping per data module type
DM_INFO_CODE_MAP: Dict[DMTypeEnum, str] = {
    DMTypeEnum.PROC: "020",
    DMTypeEnum.DESC: "010",
    DMTypeEnum.IPD: "040",
    DMTypeEnum.CIR: "050",
    DMTypeEnum.SNS: "060",
    DMTypeEnum.WIR: "070",
    DMTypeEnum.GEN: "000",
}

class DocumentService:
    """Service for document processing and management."""

    def __init__(
        self,
        upload_path: str = "/tmp/aquila_uploads",
        settings: Any | None = None,
        notifier: Callable[[str], None] | None = None,
        db: Any | None = None,
    ):
        self.upload_path = Path(upload_path)
        self.upload_path.mkdir(parents=True, exist_ok=True)
        self.icn_path = self.upload_path / "icns"
        self.icn_path.mkdir(parents=True, exist_ok=True)
        self.export_path = self.upload_path / "exports"
        self.export_path.mkdir(parents=True, exist_ok=True)
        self.settings = settings
        self.db = db
        self.notifier = notifier
        backend_root = Path(__file__).resolve().parent.parent
        self.templates_path = backend_root / "templates"
        self.schema_path = backend_root / "schemas" / "simple_data_module.xsd"
        self.audit_service = AuditService(self.upload_path / "audit.log")

    async def load_settings(self) -> Any:
        """Load settings from the database if available."""
        if self.db is not None:
            doc = await self.db.settings.find_one({})
            if doc:
                self.settings = SettingsModel(**doc)
        return self.settings

    async def upload_document(
        self,
        file_data: bytes,
        filename: str,
        mime_type: str,
        security_level: SecurityLevel = SecurityLevel.UNCLASSIFIED,
    ) -> UploadedDocument:
        """Upload and store a document."""
        sha256_hash = hashlib.sha256(file_data).hexdigest()
        file_id = str(uuid.uuid4())
        ext = Path(filename).suffix
        stored_name = f"{file_id}{ext}"
        file_path = self.upload_path / stored_name
        async with aiofiles.open(file_path, "wb") as f:
            await f.write(file_data)
        return UploadedDocument(
            filename=filename,
            file_path=str(file_path),
            mime_type=mime_type,
            file_size=len(file_data),
            sha256_hash=sha256_hash,
            security_level=security_level,
            metadata={},
        )

    async def extract_text_from_document(self, document: UploadedDocument) -> str:
        """Extract text content from a document."""
        fp = Path(document.file_path)
        try:
            if document.mime_type == "application/pdf":
                return await self._extract_pdf_text(fp)
            if (
                document.mime_type
                == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ):
                return await self._extract_docx_text(fp)
            if (
                document.mime_type
                == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ):
                return await self._extract_pptx_text(fp)
            if (
                document.mime_type
                == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
            ):
                return await self._extract_xlsx_text(fp)
            if document.mime_type.startswith("text/"):
                return await self._extract_plain_text(fp)
            return ""
        except Exception as e:
            logger.error(f"Error extracting text from {document.filename}: {e}")
            return ""

    async def _extract_pdf_text(self, file_path: Path) -> str:
        try:
            reader = PdfReader(str(file_path))
            return "".join(page.extract_text() + "\n" for page in reader.pages)
        except Exception as e:
            logger.error(f"Error extracting PDF text: {e}")
            return ""

    async def _extract_docx_text(self, file_path: Path) -> str:
        try:
            doc = await asyncio.to_thread(Document, str(file_path))
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            return "\n".join(paragraphs)
        except Exception as e:
            logger.error(f"Error extracting DOCX text: {e}")
            return ""

    async def _extract_pptx_text(self, file_path: Path) -> str:
        try:
            prs = Presentation(str(file_path))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting PPTX text: {e}")
            return ""

    async def _extract_xlsx_text(self, file_path: Path) -> str:
        try:
            workbook = load_workbook(str(file_path))
            text = ""
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            text += str(cell) + " "
                    text += "\n"
            return text
        except Exception as e:
            logger.error(f"Error extracting XLSX text: {e}")
            return ""

    async def _extract_plain_text(self, file_path: Path) -> str:
        try:
            async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
                return await f.read()
        except Exception as e:
            logger.error(f"Error extracting plain text: {e}")
            return ""

    def _derive_lcn(self, name: str) -> str:
        match = re.search(r"(LCN-[A-Za-z0-9_-]+)", name, re.IGNORECASE)
        if match:
            return match.group(1).upper()
        return f"LCN-{uuid.uuid4().hex[:8].upper()}"

    async def extract_images_from_document(
        self, document: UploadedDocument
    ) -> List[ICN]:
        if document.mime_type == "application/pdf":
            return await self._extract_pdf_images(document)
        if document.mime_type.startswith("image/"):
            return await self._process_single_image(document)
        return []

    async def _extract_pdf_images(self, document: UploadedDocument) -> List[ICN]:
        file_path = Path(document.file_path)
        icns: List[ICN] = []
        try:
            pages = await asyncio.to_thread(convert_from_path, str(file_path))
        except Exception as e:
            logger.error(f"Error extracting PDF images: {e}")
            if self.notifier:
                try:
                    self.notifier(
                        f"Failed to extract images from {document.filename}: {e}"
                    )
                except Exception:
                    pass
            return []

        for idx, img in enumerate(pages, start=1):
            filename = f"{file_path.stem}_page{idx}.png"
            out_path = self.icn_path / filename
            img.save(out_path, format="PNG")
            async with aiofiles.open(out_path, "rb") as f:
                data = await f.read()
            sha256_hash = hashlib.sha256(data).hexdigest()
            width, height = img.size
            context = pytesseract.image_to_string(img)
            icns.append(
                ICN(
                    filename=filename,
                    file_path=str(out_path),
                    sha256_hash=sha256_hash,
                    mime_type="image/png",
                    width=width,
                    height=height,
                    lcn=self._derive_lcn(filename),
                    caption=context.strip(),
                )
            )
        return icns

    async def _process_single_image(self, document: UploadedDocument) -> List[ICN]:
        try:
            async with aiofiles.open(document.file_path, "rb") as f:
                image_data = await f.read()
            image = Image.open(io.BytesIO(image_data))
            width, height = image.size
            return [
                ICN(
                    filename=document.filename,
                    file_path=document.file_path,
                    sha256_hash=document.sha256_hash,
                    mime_type=document.mime_type,
                    width=width,
                    height=height,
                    lcn=self._derive_lcn(document.filename),
                )
            ]
        except Exception as e:
            logger.error(f"Error processing image: {e}")
            return []

    def _parse_warnings_cautions(self, text: str) -> tuple[list[str], list[str]]:
        """Extract warnings and cautions from plain text."""
        warnings: list[str] = []
        cautions: list[str] = []
        for line in text.splitlines():
            stripped = line.strip()
            lower = stripped.lower()
            if lower.startswith("warning"):
                part = stripped.split(":", 1)
                warnings.append(part[1].strip() if len(part) > 1 else stripped[7:].strip())
            elif lower.startswith("caution"):
                part = stripped.split(":", 1)
                cautions.append(part[1].strip() if len(part) > 1 else stripped[7:].strip())
        return warnings, cautions

    def _format_warnings_cautions(self, warnings: list[str], cautions: list[str]) -> str:
        """Format warnings and cautions using S1000D tags."""
        parts: list[str] = []
        for w in warnings:
            parts.append(f"<warning><warningtext><para>{w}</para></warningtext></warning>")
        for c in cautions:
            parts.append(f"<caution><cautiontext><para>{c}</para></cautiontext></caution>")
        return "\n".join(parts)

    async def gather_all_documents_text(self, exclude_id: str | None = None) -> str:
        """Concatenate text from all uploaded documents."""
        if self.db is None:
            return ""
        docs = await self.db.documents.find().to_list(1000)
        texts: list[str] = []
        for d in docs:
            if exclude_id and d.get("id") == exclude_id:
                continue
            try:
                txt = await self.extract_text_from_document(UploadedDocument(**d))
                if txt:
                    texts.append(txt)
            except Exception:
                continue
        return "\n".join(texts)

    async def review_module_ai(self, content: str) -> Dict[str, Any]:
        """Use AI provider to review module content."""
        if not (os.environ.get("OPENAI_API_KEY") or os.environ.get("ANTHROPIC_API_KEY")):
            return {"issues": [], "suggested_text": content}
        provider = ProviderFactory.create_text_provider()
        req = TextProcessingRequest(text=content, task_type="review")
        res = await provider.review_module(req)
        return res.result

    async def refresh_cross_references(self) -> None:
        """Update dm_refs and icn_refs across all modules based on content."""
        if self.db is None:
            return
        modules = await self.db.data_modules.find().to_list(1000)
        icns = await self.db.icns.find().to_list(1000)
        lcn_set = {i.get("lcn") for i in icns}
        dmc_set = {m.get("dmc") for m in modules}
        for m in modules:
            dm_refs = set(m.get("dm_refs", []))
            icn_refs = set(m.get("icn_refs", []))
            content = m.get("content", "")
            for dmc in dmc_set:
                if dmc != m.get("dmc") and dmc in content:
                    dm_refs.add(dmc)
            for lcn in lcn_set:
                if lcn in content:
                    icn_refs.add(lcn)
            if dm_refs != set(m.get("dm_refs", [])) or icn_refs != set(m.get("icn_refs", [])):
                await self.db.data_modules.update_one(
                    {"dmc": m.get("dmc")},
                    {"$set": {"dm_refs": list(dm_refs), "icn_refs": list(icn_refs), "updated_at": datetime.utcnow()}}
                )

    async def process_document_with_ai(
        self, document: UploadedDocument, text_content: str
    ) -> List[DataModule]:
        await self.load_settings()
        text_provider = ProviderFactory.create_text_provider()
        if self.db is not None:
            extra_text = await self.gather_all_documents_text(exclude_id=document.id)
            if extra_text:
                text_content = f"{text_content}\n{extra_text}"[:10000]
        try:
            classification = TextProcessingRequest(
                text=text_content, task_type="classify"
            )
            class_response = await text_provider.classify_document(classification)
            if "error" in class_response.result:
                raise Exception(class_response.result["error"])

            extract_req = TextProcessingRequest(text=text_content, task_type="extract")
            extract_res = await text_provider.extract_structured_data(extract_req)
            if "error" in extract_res.result:
                raise Exception(extract_res.result["error"])

            refs = extract_res.result.get("references", [])
            dm_refs = [r["reference"] for r in refs if r.get("type") == "dm"]
            icn_refs = [
                self._derive_lcn(r["reference"])
                for r in refs
                if r.get("type") in {"figure", "image", "table"}
            ]

            warn_provider = extract_res.result.get("warnings", [])
            caution_provider = extract_res.result.get("cautions", [])
            warn_text, caution_text = self._parse_warnings_cautions(text_content)
            warnings = list({*warn_provider, *warn_text})
            cautions = list({*caution_provider, *caution_text})
            wc_prefix = self._format_warnings_cautions(warnings, cautions)

            verbatim = DataModule(
                dmc=self._generate_dmc(class_response.result),
                title=class_response.result.get("title", "Untitled Document"),
                dm_type=DMTypeEnum(class_response.result.get("dm_type", "GEN")),
                info_variant="00",
                content="{}\n{}".format(wc_prefix, text_content).strip(),
                source_document_id=document.id,
                security_level=document.security_level,
                processing_status="completed",
                dm_refs=dm_refs,
                icn_refs=icn_refs,
            )

            rewrite_req = TextProcessingRequest(text=text_content, task_type="rewrite")
            rewrite_res = await text_provider.rewrite_to_ste(rewrite_req)

            modules = [verbatim]
            if "error" not in rewrite_res.result:
                ste_dm = DataModule(
                    dmc=self._generate_dmc(class_response.result, variant="01"),
                    title=class_response.result.get("title", "Untitled Document"),
                    dm_type=DMTypeEnum(class_response.result.get("dm_type", "GEN")),
                    info_variant="01",
                    content="{}\n{}".format(
                        wc_prefix,
                        rewrite_res.result.get("rewritten_text", text_content),
                    ).strip(),
                    source_document_id=document.id,
                    security_level=document.security_level,
                    ste_score=rewrite_res.result.get("ste_score", 0.0),
                    processing_status="completed",
                    dm_refs=dm_refs,
                    icn_refs=icn_refs,
                )
                modules.append(ste_dm)
            return modules
        except Exception as e:
            logger.error(f"Error processing document with AI: {e}")
            basic = DataModule(
                dmc=self._generate_dmc({"dm_type": "GEN", "title": document.filename}),
                title=document.filename,
                dm_type=DMTypeEnum.GEN,
                info_variant="00",
                content="{}\n{}".format(self._format_warnings_cautions(*self._parse_warnings_cautions(text_content)), text_content).strip(),
                source_document_id=document.id,
                security_level=document.security_level,
                processing_status="error",
                dm_refs=[],
                icn_refs=[],
            )
            return [basic]

    def _generate_dmc(self, classification_result: dict, variant: str = "00") -> str:
        """Generate a fully S1000D compliant Data Module Code."""
        cfg = self.settings.dmc_defaults if self.settings else {}
        structure = DEFAULT_STRUCTURE_CODES.get(
            getattr(self.settings, "structure_type", StructureType.OTHER),
            DEFAULT_STRUCTURE_CODES[StructureType.OTHER],
        )

        model_ident = cfg.get("model_ident", "AQUILA")
        system_diff = structure.get("system_diff", cfg.get("system_diff", "00"))
        system_code = structure.get("system_code", cfg.get("system_code", "000"))
        sub_system_code = structure.get("sub_system_code", cfg.get("sub_system_code", "00"))
        sub_sub_system_code = structure.get("sub_sub_system_code", cfg.get("sub_sub_system_code", "00"))

        assy_code = cfg.get("assy_code", "00")
        disassy_code = cfg.get("disassy_code", "00")
        disassy_code_variant = cfg.get("disassy_code_variant", "00")

        dm_type = DMTypeEnum(classification_result.get("dm_type", "GEN"))
        info_code = DM_INFO_CODE_MAP.get(dm_type, cfg.get("info_code", "000"))
        info_code_variant = cfg.get("info_code_variant", "A")
        item_location_code = cfg.get("item_location_code", "A")
        learn_code = cfg.get("learn_code", "00")
        learn_event_code = cfg.get("learn_event_code", "00")

        return (
            f"DMC-{model_ident}-{system_diff}-{system_code}-{sub_system_code}-"
            f"{sub_sub_system_code}-{assy_code}-{disassy_code}-{disassy_code_variant}-"
            f"{info_code}-{info_code_variant}-{item_location_code}-{learn_code}-"
            f"{learn_event_code}-{variant}"
        )

    async def process_image_with_ai(self, icn: ICN) -> ICN:
        vision_provider = ProviderFactory.create_vision_provider()
        try:
            async with aiofiles.open(icn.file_path, "rb") as f:
                image_data = await f.read()
            image_base64 = base64.b64encode(image_data).decode("utf-8")
            caption_req = VisionProcessingRequest(
                image_data=image_base64, task_type="caption"
            )
            caption_res = await vision_provider.generate_caption(caption_req)
            objects_req = VisionProcessingRequest(
                image_data=image_base64, task_type="objects"
            )
            objects_res = await vision_provider.detect_objects(objects_req)
            hotspots_req = VisionProcessingRequest(
                image_data=image_base64, task_type="hotspots"
            )
            hotspots_res = await vision_provider.generate_hotspots(hotspots_req)
            icn.caption = caption_res.caption
            icn.objects = objects_res.objects
            icn.hotspots = hotspots_res.hotspots
            return icn
        except Exception as e:
            logger.error(f"Error processing image with AI: {e}")
            icn.caption = f"Error processing image: {e}"
            return icn

    def _render_pdf(self, module: DataModule, icns: List[ICN], pdf_path: Path) -> None:
        """Render a PDF file for the given data module."""
        styles = getSampleStyleSheet()
        styles.add(
            ParagraphStyle(
                name="Caption",
                parent=styles["Normal"],
                fontSize=10,
                italic=True,
            )
        )

        flow: List[Any] = [Paragraph(module.title, styles["Title"]), Spacer(1, 12)]

        for para in module.content.split("\n\n"):
            txt = para.strip()
            if txt:
                flow.append(Paragraph(txt, styles["BodyText"]))
                flow.append(Spacer(1, 12))

        for icn in icns:
            try:
                flow.append(Image(icn.file_path, width=400, preserveAspectRatio=True))
                if icn.caption:
                    flow.append(Paragraph(icn.caption, styles["Caption"]))
                flow.append(Spacer(1, 12))
            except Exception as e:  # pragma: no cover - best effort logging
                logger.warning(f"Could not add image {icn.file_path}: {e}")

        doc = SimpleDocTemplate(str(pdf_path))
        doc.build(flow)

    def render_data_module_xml(self, module: DataModule) -> str:
        """Render a DataModule to XML using Jinja2 template."""
        env = Environment(
            loader=FileSystemLoader(str(self.templates_path)),
            autoescape=select_autoescape(["xml"]),
        )
        template = env.get_template("data_module.xml.j2")
        return template.render(module=module)

    def validate_xml(self, xml_str: str) -> bool:
        """Validate XML string against built-in XSD."""
        try:
            schema = xmlschema.XMLSchema(self.schema_path)
            return schema.is_valid(xml_str)
        except Exception:
            return False

    async def publish_publication_module(
        self,
        pm: PublicationModule,
        db,
        formats: List[str],
        variants: List[str],
    ) -> Dict[str, Any]:
        """Compile a publication module into an export package."""
        modules = await db.data_modules.find({"dmc": {"$in": pm.dm_list}}).to_list(
            len(pm.dm_list)
        )
        if not modules:
            raise ValueError("No data modules found for publication")

        pm_dir = self.export_path / pm.pm_code
        pm_dir.mkdir(parents=True, exist_ok=True)

        package_files: List[Path] = []
        errors: List[str] = []
        env = Environment(loader=FileSystemLoader(str(self.templates_path)))
        html_template = env.get_template("data_module.html.j2")

        for mod_data in modules:
            dm = DataModule(**mod_data)
            if dm.info_variant not in variants:
                continue

            try:
                xml_str = self.render_data_module_xml(dm)
                xml_path = pm_dir / f"{dm.dmc}_{dm.info_variant}.xml"
                async with aiofiles.open(xml_path, "w") as f:
                    await f.write(xml_str)
                package_files.append(xml_path)

                if "html" in formats:
                    html_path = pm_dir / f"{dm.dmc}_{dm.info_variant}.html"
                    html_content = html_template.render(module=dm)
                    async with aiofiles.open(html_path, "w") as f:
                        await f.write(html_content)
                    package_files.append(html_path)

                if "pdf" in formats:
                    pdf_path = pm_dir / f"{dm.dmc}_{dm.info_variant}.pdf"
                    icn_objs: List[ICN] = []
                    if dm.icn_refs:
                        try:
                            icn_data = await db.icns.find({"lcn": {"$in": dm.icn_refs}}).to_list(len(dm.icn_refs))
                            icn_objs = [ICN(**i) for i in icn_data]
                        except Exception as e:  # pragma: no cover - best effort logging
                            logger.warning(f"Failed to load ICNs for {dm.dmc}: {e}")
                    await asyncio.to_thread(self._render_pdf, dm, icn_objs, pdf_path)
                    package_files.append(pdf_path)
            except Exception as e:  # pragma: no cover - best effort logging
                error_msg = f"Failed to export {dm.dmc}: {e}"
                logger.warning(error_msg)
                errors.append(error_msg)

        if not package_files:
            raise ValueError("No files generated for publication module")

        zip_path = shutil.make_archive(str(pm_dir), "zip", pm_dir)
        return {"package": Path(zip_path), "errors": errors}
